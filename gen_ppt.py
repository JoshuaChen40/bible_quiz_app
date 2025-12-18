import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- 1. 配置與顏色定義 ---
MAIN_TITLE = "✨ 腦光一閃"
LIST_PAGE_TITLE = "📚 題目清單索引"

COLOR_HIGHLIGHT = RGBColor(0x28, 0x7E, 0xF3)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUCCESS = RGBColor(0x00, 0xAA, 0x00)
BUFFER_IMAGE_PATH = "your_buffer_image.png"

GROUP_COLORS = [0x287EF3, 0x28B463, 0xD4AC0D, 0xCB4335, 0x884EA0, 0x17A589, 0xD35400]

# --- 2. 輔助函數 ---

def add_nav_back_button(slide):
    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(4.7), Inches(2.4), Inches(0.7)
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = RGBColor(200, 200, 200)
    btn.line.color.rgb = COLOR_WHITE
    tf = btn.text_frame
    tf.text = "🏠 回列表"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(50, 50, 50)
    btn.click_action.hyperlink.address = "#2"

def create_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    return slide

def add_grouped_index_to_slide(prs, questions, list_title):
    """ 集合頁修正：確保方塊與文字的超連結皆正確指向 #{3 + i*3} """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
    title_box.text_frame.text = list_title
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True

    unique_groups = list(dict.fromkeys([q.get('q_group', '一般') for q in questions]))
    group_color_map = {name: GROUP_COLORS[i % len(GROUP_COLORS)] for i, name in enumerate(unique_groups)}

    cols = 6
    width, height = Inches(1.45), Inches(0.65)
    margin_x, margin_y = Inches(0.12), Inches(0.1)
    start_x, start_y = Inches(0.3), Inches(0.7)

    for i, q in enumerate(questions):
        if i >= 36: break
        row, col = i // cols, i % cols
        left, top = start_x + (col * (width + margin_x)), start_y + (row * (height + margin_y))
        
        # 1. 建立目標索引
        target_q_idx = 3 + (i * 3)
        target_link = f"#{target_q_idx}"

        # 2. 建立方塊並設定方塊點擊連結
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        current_color = group_color_map.get(q.get('q_group', '一般'))
        shape.fill.fore_color.rgb = RGBColor.from_string(f"{current_color:06x}")
        shape.line.color.rgb = COLOR_WHITE
        shape.click_action.hyperlink.address = target_link # 🎯 方塊點擊修正

        # 3. 建立文字並設定文字連結 (為了支援點擊後變色)
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = 1
        
        run = p.add_run()
        run.text = f"Q{i+1}\n{q.get('q_group', '')[:5]}"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE
        run.hyperlink.address = target_link # 🎯 文字連結修正

def add_blank_buffer_slide(prs, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if image_path and os.path.exists(image_path):
        try: slide.shapes.add_picture(image_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        except: pass
    add_nav_back_button(slide)
    return slide

def add_question_slide(prs, q_idx, q_data, target_ans_idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(7), Inches(0.4))
    title.text_frame.text = f"第 {q_idx + 1} 題 | {q_data.get('q_group', '')}"
    title.text_frame.paragraphs[0].font.size = Pt(20)

    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(1.8))
    q_box.text_frame.word_wrap = True
    p_q = q_box.text_frame.paragraphs[0]
    p_q.text = q_data["question"]
    p_q.font.size = Pt(38)
    p_q.font.bold = True

    opt_pos = {"A": (0.5, 2.7), "B": (5.2, 2.7), "C": (0.5, 3.8), "D": (5.2, 3.8)}
    for opt, (l, t) in opt_pos.items():
        if opt in q_data and q_data[opt]:
            ob = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(4.5), Inches(0.8))
            p = ob.text_frame.paragraphs[0]
            p.text = f"({opt}) {q_data[opt]}"
            p.font.size = Pt(35)
            p.font.bold = True

    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(0.15), Inches(1.8), Inches(0.5))
    btn.fill.solid()
    btn.fill.fore_color.rgb = COLOR_HIGHLIGHT
    btn.text_frame.text = "📜 看解答"
    btn.click_action.hyperlink.address = f"#{target_ans_idx}"
    return slide

def add_answer_slide(prs, q_idx, q_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ans_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
    ans_box.text_frame.text = f"正確答案：{q_data['answer']}"
    p_ans = ans_box.text_frame.paragraphs[0]
    p_ans.font.size = Pt(58); p_ans.font.bold = True; p_ans.font.color.rgb = COLOR_SUCCESS

    exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(3.5))
    exp_box.text_frame.word_wrap = True
    tf = exp_box.text_frame
    p_label = tf.paragraphs[0]
    p_label.text = "💡 解釋："
    p_label.font.size = Pt(25)
    p_content = tf.add_paragraph()
    p_content.text = q_data.get('explanation', '無')
    p_content.font.size = Pt(55)
    p_content.font.bold = True
    
    add_nav_back_button(slide)
    return slide

# --- 3. 主流程 ---

def generate_quiz_pptx_final(json_path, output_name="Quiz.pptx"):
    if not os.path.exists(json_path): return
    with open(json_path, "r", encoding="utf-8") as f:
        questions = json.load(f)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
    
    create_title_slide(prs, MAIN_TITLE, "點擊方塊進入題目")
    add_grouped_index_to_slide(prs, questions, LIST_PAGE_TITLE) 

    ans_start = 2 + (len(questions) * 3) + 1

    for i, q in enumerate(questions):
        # 題目頁索引：3, 6, 9...
        target_ans = ans_start + (i * 3)
        add_question_slide(prs, i, q, target_ans)
        add_blank_buffer_slide(prs, BUFFER_IMAGE_PATH)
        add_blank_buffer_slide(prs, BUFFER_IMAGE_PATH)

    for i, q in enumerate(questions):
        add_answer_slide(prs, i, q)
        add_blank_buffer_slide(prs, BUFFER_IMAGE_PATH)
        add_blank_buffer_slide(prs, BUFFER_IMAGE_PATH)

    prs.save(output_name)
    print(f"✅ 生成成功：{output_name}")

if __name__ == "__main__":
    generate_quiz_pptx_final("questions_data.json", f"{MAIN_TITLE}_互動簡報.pptx")